from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Create presentation
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Barriers to Effective Visual Communication"
    subtitle.text = "Understanding and Overcoming Common Pitfalls\nGenerated via Python-PPTX"

    # Slide 1: Introduction
    bullet_slide_layout = prs.slide_layouts[1]
    slide1 = prs.slides.add_slide(bullet_slide_layout)
    shapes1 = slide1.shapes
    title_shape1 = shapes1.title
    body_shape1 = shapes1.placeholders[1]

    title_shape1.text = "Introduction"
    tf1 = body_shape1.text_frame
    tf1.text = "Visual communication is powerful, but it can fail if barriers exist."

    p1 = tf1.add_paragraph()
    p1.text = "Common barriers prevent the audience from understanding the intended message."
    p1.level = 1

    p2 = tf1.add_paragraph()
    p2.text = "Key barriers include:"
    p2.level = 1

    p3 = tf1.add_paragraph()
    p3.text = "Information Overload"
    p3.level = 2

    p4 = tf1.add_paragraph()
    p4.text = "Poor Design and Clutter"
    p4.level = 2

    p5 = tf1.add_paragraph()
    p5.text = "Inappropriate Color Choices"
    p5.level = 2

    p6 = tf1.add_paragraph()
    p6.text = "Cultural Differences"
    p6.level = 2

    p7 = tf1.add_paragraph()
    p7.text = "Misleading Data Visualizations"
    p7.level = 2

    # Slide 2: Information Overload
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes2 = slide2.shapes
    title_shape2 = shapes2.title
    body_shape2 = shapes2.placeholders[1]

    title_shape2.text = "Barrier 1: Information Overload"
    tf2 = body_shape2.text_frame
    tf2.text = "Too much information on a single visual can overwhelm the audience."

    p = tf2.add_paragraph()
    p.text = "Too much text or data makes it difficult to focus on key points."
    p.level = 1

    p = tf2.add_paragraph()
    p.text = "Solution: Keep it simple."
    p.level = 1

    p = tf2.add_paragraph()
    p.text = "Use whitespace effectively."
    p.level = 2

    p = tf2.add_paragraph()
    p.text = "Stick to one main idea per slide or visual."
    p.level = 2

    # Slide 3: Poor Design and Clutter
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    shapes3 = slide3.shapes
    title_shape3 = shapes3.title
    body_shape3 = shapes3.placeholders[1]

    title_shape3.text = "Barrier 2: Poor Design and Clutter"
    tf3 = body_shape3.text_frame
    tf3.text = "A lack of visual hierarchy and cluttered elements confuse viewers."

    p = tf3.add_paragraph()
    p.text = "Inconsistent fonts and unaligned elements look unprofessional."
    p.level = 1

    p = tf3.add_paragraph()
    p.text = "Unnecessary graphics distract from the core message."
    p.level = 1

    p = tf3.add_paragraph()
    p.text = "Solution: Use grids, consistent typography, and remove non-essential elements."
    p.level = 1

    # Slide 4: Inappropriate Color Choices
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    shapes4 = slide4.shapes
    title_shape4 = shapes4.title
    body_shape4 = shapes4.placeholders[1]

    title_shape4.text = "Barrier 3: Inappropriate Color Choices"
    tf4 = body_shape4.text_frame
    tf4.text = "Poor contrast and relying only on color can alienate parts of the audience."

    p = tf4.add_paragraph()
    p.text = "Low contrast makes text hard to read."
    p.level = 1

    p = tf4.add_paragraph()
    p.text = "Color blindness affects approximately 8% of men and 0.5% of women."
    p.level = 1

    p = tf4.add_paragraph()
    p.text = "Solution: Ensure high contrast and use patterns/labels alongside color."
    p.level = 1

    # Slide 5: Cultural Differences
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    shapes5 = slide5.shapes
    title_shape5 = shapes5.title
    body_shape5 = shapes5.placeholders[1]

    title_shape5.text = "Barrier 4: Cultural Differences"
    tf5 = body_shape5.text_frame
    tf5.text = "Symbols, colors, and imagery can mean different things across cultures."

    p = tf5.add_paragraph()
    p.text = "Example: Red can mean danger/stop in Western cultures, but prosperity in some Eastern cultures."
    p.level = 1

    p = tf5.add_paragraph()
    p.text = "Reading direction (left-to-right vs. right-to-left) affects visual flow."
    p.level = 1

    p = tf5.add_paragraph()
    p.text = "Solution: Know your audience and localize visual content when necessary."
    p.level = 1

    # Slide 6: Misleading Data Visualizations
    slide6 = prs.slides.add_slide(bullet_slide_layout)
    shapes6 = slide6.shapes
    title_shape6 = shapes6.title
    body_shape6 = shapes6.placeholders[1]

    title_shape6.text = "Barrier 5: Misleading Data Visualizations"
    tf6 = body_shape6.text_frame
    tf6.text = "Charts and graphs that misrepresent data cause loss of trust."

    p = tf6.add_paragraph()
    p.text = "Truncated y-axes can exaggerate small differences."
    p.level = 1

    p = tf6.add_paragraph()
    p.text = "Improper chart types (e.g., using a pie chart for too many categories)."
    p.level = 1

    p = tf6.add_paragraph()
    p.text = "Solution: Follow data visualization best practices and maintain integrity."
    p.level = 1

    # Slide 7: Conclusion
    slide7 = prs.slides.add_slide(bullet_slide_layout)
    shapes7 = slide7.shapes
    title_shape7 = shapes7.title
    body_shape7 = shapes7.placeholders[1]

    title_shape7.text = "Conclusion"
    tf7 = body_shape7.text_frame
    tf7.text = "Effective visual communication requires intentional design."

    p = tf7.add_paragraph()
    p.text = "Avoid overload, clutter, and misleading representations."
    p.level = 1

    p = tf7.add_paragraph()
    p.text = "Consider accessibility (color) and audience background (culture)."
    p.level = 1

    p = tf7.add_paragraph()
    p.text = "Clear, simple, and accurate visuals lead to successful communication."
    p.level = 1

    # Save presentation
    prs.save("barriers_of_effective_visual_communication.pptx")
    print("Presentation saved successfully as barriers_of_effective_visual_communication.pptx")

if __name__ == "__main__":
    create_presentation()
